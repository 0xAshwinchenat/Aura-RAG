import io
import re
import csv
import zipfile
import logging
from typing import List, Dict, Any
from pydantic import BaseModel
import filetype

# Import parsers
import pypdf
import docx
import pptx
import openpyxl
from bs4 import BeautifulSoup
from email import message_from_bytes
from PIL import Image

try:
    import pytesseract
except ImportError:
    pytesseract = None

logger = logging.getLogger(__name__)

class ParsedSection(BaseModel):
    text: str
    metadata: Dict[str, Any]

class ParsedDocument(BaseModel):
    source_name: str
    mime_type: str
    sections: List[ParsedSection]
    success: bool
    error_message: str | None = None

def detect_mime_type(file_path: str) -> str:
    """
    Detects the file MIME type reliably by checking the file signatures (magic numbers)
    and zip file contents, rather than trusting the extension.
    """
    # 1. Try binary signature detection using filetype library
    kind = filetype.guess(file_path)
    if kind:
        mime = kind.mime
        
        # If it's a zip file, it might be a DOCX, PPTX, or XLSX
        if mime == "application/zip":
            try:
                with zipfile.ZipFile(file_path, 'r') as zip_ref:
                    namelist = zip_ref.namelist()
                    if "word/document.xml" in namelist:
                        return "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                    elif "xl/workbook.xml" in namelist:
                        return "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                    elif "ppt/presentation.xml" in namelist:
                        return "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            except Exception:
                pass
        return mime

    # 2. If binary check yields nothing, read as text and analyze signatures
    try:
        with open(file_path, 'rb') as f:
            header_bytes = f.read(2048)
        
        # Check standard PDF signature (in case filetype missed it)
        if header_bytes.startswith(b"%PDF-"):
            return "application/pdf"
            
        # Try decoding as text
        for encoding in ['utf-8', 'latin-1', 'utf-16']:
            try:
                header_text = header_bytes.decode(encoding)
                
                # Check for HTML
                if re.search(r'<!DOCTYPE\s+html|<html|<head|<body', header_text, re.IGNORECASE):
                    return "text/html"
                
                # Check for Email (EML) - common headers
                eml_headers = ["from:", "to:", "subject:", "date:", "received:", "mime-version:", "delivered-to:"]
                lines = [line.strip().lower() for line in header_text.splitlines() if line.strip()][:10]
                if any(any(line.startswith(h) for h in eml_headers) for line in lines):
                    return "message/rfc822"
                
                # Check for CSV
                try:
                    # Sniff CSV behavior
                    dialect = csv.Sniffer().sniff(header_text[:1000])
                    if dialect.delimiter in [',', ';', '\t']:
                        # double check columns on first couple of lines
                        rows = list(csv.reader(header_text.splitlines()[:3], dialect))
                        if len(rows) > 1 and len(rows[0]) == len(rows[1]) and len(rows[0]) > 1:
                            return "text/csv"
                except Exception:
                    pass
                
                # Check for Markdown
                if re.search(r'^#\s+|^##\s+|^###\s+|\*\*.*?\*\*|\[.*?\]\(.*?\)', header_text, re.MULTILINE):
                    return "text/markdown"
                
                # Fallback to plain text if it successfully decoded
                return "text/plain"
            except UnicodeDecodeError:
                continue
    except Exception as e:
        logger.error(f"Error sniffing file bytes for {file_path}: {e}")

    # Ultimate fallback
    return "application/octet-stream"


class DocumentParser:
    @staticmethod
    def parse_pdf(file_path: str, source_name: str) -> ParsedDocument:
        sections = []
        try:
            reader = pypdf.PdfReader(file_path)
            
            # Check for encryption/password protection
            if reader.is_encrypted:
                return ParsedDocument(
                    source_name=source_name,
                    mime_type="application/pdf",
                    sections=[],
                    success=False,
                    error_message="File is password-protected or encrypted."
                )
            
            num_pages = len(reader.pages)
            if num_pages == 0:
                return ParsedDocument(
                    source_name=source_name,
                    mime_type="application/pdf",
                    sections=[],
                    success=False,
                    error_message="PDF file is empty."
                )

            for i, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                clean_text = text.strip()
                
                # Fallback to OCR if page text is empty and pytesseract is available
                if not clean_text and pytesseract is not None:
                    # PyPDF allows extracting images, or we could run OCR on a rendered page.
                    # Since we want to fail gracefully, we will log a warning.
                    # Rendered PDF OCR is heavy, but let's check for basic OCR if possible.
                    pass
                
                sections.append(ParsedSection(
                    text=clean_text,
                    metadata={"page": i + 1, "source": source_name}
                ))
            
            # Check if all pages are completely empty
            total_text = "".join([s.text for s in sections]).strip()
            if not total_text:
                return ParsedDocument(
                    source_name=source_name,
                    mime_type="application/pdf",
                    sections=sections,
                    success=True,
                    error_message="Parsed successfully, but no text could be extracted (possible scanned PDF or images only)."
                )

            return ParsedDocument(
                source_name=source_name,
                mime_type="application/pdf",
                sections=sections,
                success=True
            )
        except Exception as e:
            return ParsedDocument(
                source_name=source_name,
                mime_type="application/pdf",
                sections=[],
                success=False,
                error_message=f"PyPDF error: {str(e)}"
            )

    @staticmethod
    def parse_docx(file_path: str, source_name: str) -> ParsedDocument:
        try:
            doc = docx.Document(file_path)
            sections = []
            current_section_name = "Header/Introduction"
            current_paragraphs = []
            
            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if not text:
                    continue
                
                # Detect structural headings to group content into sections
                if paragraph.style and paragraph.style.name and paragraph.style.name.startswith("Heading"):
                    if current_paragraphs:
                        sections.append(ParsedSection(
                            text="\n".join(current_paragraphs),
                            metadata={"section": current_section_name, "source": source_name}
                        ))
                        current_paragraphs = []
                    current_section_name = text
                else:
                    current_paragraphs.append(text)
            
            # Handle tables
            table_text_list = []
            for i, table in enumerate(doc.tables):
                table_rows = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_rows.append(" | ".join(row_data))
                if table_rows:
                    table_text_list.append(f"Table {i+1}:\n" + "\n".join(table_rows))
            
            if current_paragraphs:
                sections.append(ParsedSection(
                    text="\n".join(current_paragraphs),
                    metadata={"section": current_section_name, "source": source_name}
                ))
                
            if table_text_list:
                sections.append(ParsedSection(
                    text="\n\n".join(table_text_list),
                    metadata={"section": "Tables", "source": source_name}
                ))

            if not sections:
                return ParsedDocument(
                    source_name=source_name,
                    mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    sections=[],
                    success=False,
                    error_message="Document is empty."
                )

            return ParsedDocument(
                source_name=source_name,
                mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                sections=sections,
                success=True
            )
        except Exception as e:
            return ParsedDocument(
                source_name=source_name,
                mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                sections=[],
                success=False,
                error_message=f"python-docx error: {str(e)}"
            )

    @staticmethod
    def parse_pptx(file_path: str, source_name: str) -> ParsedDocument:
        try:
            prs = pptx.Presentation(file_path)
            sections = []
            
            for i, slide in enumerate(prs.slides):
                slide_texts = []
                slide_title = f"Slide {i+1}"
                
                # Check for shapes with text
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                slide_texts.append(text)
                                # Try to capture the main title of the slide
                                if shape.name.startswith("Title") or shape.name.startswith("Centered Title"):
                                    slide_title = f"Slide {i+1}: {text}"
                
                if slide_texts:
                    sections.append(ParsedSection(
                        text="\n".join(slide_texts),
                        metadata={"page": i + 1, "section": slide_title, "source": source_name}
                    ))
            
            if not sections:
                return ParsedDocument(
                    source_name=source_name,
                    mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    sections=[],
                    success=False,
                    error_message="Presentation contains no text or is empty."
                )

            return ParsedDocument(
                source_name=source_name,
                mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                sections=sections,
                success=True
            )
        except Exception as e:
            return ParsedDocument(
                source_name=source_name,
                mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                sections=[],
                success=False,
                error_message=f"python-pptx error: {str(e)}"
            )

    @staticmethod
    def parse_xlsx(file_path: str, source_name: str) -> ParsedDocument:
        try:
            wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
            sections = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                rows = []
                
                # Load a max of 1000 rows to prevent huge Excel files from crashing RAM
                for r_idx, row in enumerate(sheet.iter_rows(values_only=True)):
                    if r_idx > 1000:
                        rows.append(["... [Truncated due to size] ..."])
                        break
                    
                    # Convert cells to string, filtering out empty rows
                    row_values = [str(cell).strip() if cell is not None else "" for cell in row]
                    if any(val != "" for val in row_values):
                        rows.append(row_values)
                
                if rows:
                    # Convert to Markdown Table format
                    md_table = []
                    # Header row
                    md_table.append("| " + " | ".join(rows[0]) + " |")
                    # Separator
                    md_table.append("| " + " | ".join(["---"] * len(rows[0])) + " |")
                    # Data rows
                    for row in rows[1:]:
                        # Align length in case of row length mismatch
                        if len(row) < len(rows[0]):
                            row = row + [""] * (len(rows[0]) - len(row))
                        elif len(row) > len(rows[0]):
                            row = row[:len(rows[0])]
                        md_table.append("| " + " | ".join(row) + " |")
                        
                    sections.append(ParsedSection(
                        text="\n".join(md_table),
                        metadata={"sheet": sheet_name, "source": source_name}
                    ))
            
            wb.close()
            
            if not sections:
                return ParsedDocument(
                    source_name=source_name,
                    mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    sections=[],
                    success=False,
                    error_message="Excel workbook is empty."
                )

            return ParsedDocument(
                source_name=source_name,
                mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                sections=sections,
                success=True
            )
        except Exception as e:
            return ParsedDocument(
                source_name=source_name,
                mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                sections=[],
                success=False,
                error_message=f"openpyxl error: {str(e)}"
            )

    @staticmethod
    def parse_csv(file_path: str, source_name: str) -> ParsedDocument:
        try:
            # Detect encoding
            encoding = 'utf-8'
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    f.read(100)
            except UnicodeDecodeError:
                encoding = 'latin-1'
                
            with open(file_path, 'r', encoding=encoding) as f:
                # Sniff CSV details
                content = f.read(5000)
                f.seek(0)
                try:
                    dialect = csv.Sniffer().sniff(content)
                except Exception:
                    dialect = csv.excel # fallback
                
                reader = csv.reader(f, dialect)
                rows = list(reader)
                
            if not rows:
                return ParsedDocument(
                    source_name=source_name,
                    mime_type="text/csv",
                    sections=[],
                    success=False,
                    error_message="CSV is empty."
                )
            
            # Format as Markdown Table
            md_table = []
            headers = rows[0]
            md_table.append("| " + " | ".join(headers) + " |")
            md_table.append("| " + " | ".join(["---"] * len(headers)) + " |")
            for row in rows[1:]:
                # Align columns
                if len(row) < len(headers):
                    row = row + [""] * (len(headers) - len(row))
                elif len(row) > len(headers):
                    row = row[:len(headers)]
                md_table.append("| " + " | ".join([r.strip() for r in row]) + " |")
                
            return ParsedDocument(
                source_name=source_name,
                mime_type="text/csv",
                sections=[ParsedSection(text="\n".join(md_table), metadata={"section": "Table", "source": source_name})],
                success=True
            )
        except Exception as e:
            return ParsedDocument(
                source_name=source_name,
                mime_type="text/csv",
                sections=[],
                success=False,
                error_message=f"CSV parser error: {str(e)}"
            )

    @staticmethod
    def parse_html(file_path: str, source_name: str) -> ParsedDocument:
        try:
            encoding = 'utf-8'
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
            except UnicodeDecodeError:
                encoding = 'latin-1'
                with open(file_path, 'r', encoding=encoding) as f:
                    content = f.read()
                    
            soup = BeautifulSoup(content, 'html.parser')
            
            # Remove scripts and styling elements
            for script in soup(["script", "style"]):
                script.extract()
                
            # Document structure parsing (heading based)
            sections = []
            title = soup.title.string.strip() if soup.title else "HTML Document"
            
            current_section = title
            current_texts = []
            
            # Traverse HTML body elements
            body = soup.body if soup.body else soup
            for tag in body.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'li', 'td', 'pre']):
                txt = tag.get_text().strip()
                if not txt:
                    continue
                if tag.name.startswith('h'):
                    if current_texts:
                        sections.append(ParsedSection(
                            text="\n".join(current_texts),
                            metadata={"section": current_section, "source": source_name}
                        ))
                        current_texts = []
                    current_section = txt
                else:
                    current_texts.append(txt)
                    
            if current_texts:
                sections.append(ParsedSection(
                    text="\n".join(current_texts),
                    metadata={"section": current_section, "source": source_name}
                ))

                
            # Fallback if no hierarchical sections found
            if not sections:
                text = soup.get_text()
                # Clean up lines
                lines = (line.strip() for line in text.splitlines())
                chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
                clean_text = "\n".join(chunk for chunk in chunks if chunk)
                
                if clean_text:
                    sections.append(ParsedSection(
                        text=clean_text,
                        metadata={"section": title, "source": source_name}
                    ))
            
            if not sections:
                return ParsedDocument(
                    source_name=source_name,
                    mime_type="text/html",
                    sections=[],
                    success=False,
                    error_message="HTML file contains no textual content."
                )

            return ParsedDocument(
                source_name=source_name,
                mime_type="text/html",
                sections=sections,
                success=True
            )
        except Exception as e:
            return ParsedDocument(
                source_name=source_name,
                mime_type="text/html",
                sections=[],
                success=False,
                error_message=f"HTML parser error: {str(e)}"
            )

    @staticmethod
    def parse_markdown(file_path: str, source_name: str) -> ParsedDocument:
        try:
            # We can parse Markdown sections based on '#' characters
            encoding = 'utf-8'
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
            except UnicodeDecodeError:
                encoding = 'latin-1'
                with open(file_path, 'r', encoding=encoding) as f:
                    content = f.read()

            sections = []
            current_section = "Header"
            current_texts = []
            
            for line in content.splitlines():
                stripped = line.strip()
                if stripped.startswith('#'):
                    # Found a header
                    if current_texts:
                        sections.append(ParsedSection(
                            text="\n".join(current_texts),
                            metadata={"section": current_section, "source": source_name}
                        ))
                        current_texts = []
                    # Clean up header markup
                    current_section = stripped.lstrip('#').strip()
                else:
                    if stripped:
                        current_texts.append(stripped)
                        
            if current_texts:
                sections.append(ParsedSection(
                    text="\n".join(current_texts),
                    metadata={"section": current_section, "source": source_name}
                ))
                
            if not sections:
                return ParsedDocument(
                    source_name=source_name,
                    mime_type="text/markdown",
                    sections=[],
                    success=False,
                    error_message="Markdown file is empty."
                )

            return ParsedDocument(
                source_name=source_name,
                mime_type="text/markdown",
                sections=sections,
                success=True
            )
        except Exception as e:
            return ParsedDocument(
                source_name=source_name,
                mime_type="text/markdown",
                sections=[],
                success=False,
                error_message=f"Markdown parser error: {str(e)}"
            )

    @staticmethod
    def parse_eml(file_path: str, source_name: str) -> ParsedDocument:
        try:
            with open(file_path, 'rb') as f:
                msg = message_from_bytes(f.read())
            
            headers = {
                "from": msg.get("From", "Unknown Sender"),
                "to": msg.get("To", "Unknown Recipient"),
                "subject": msg.get("Subject", "No Subject"),
                "date": msg.get("Date", "No Date")
            }
            
            body_text = ""
            if msg.is_multipart():
                for part in msg.walk():
                    content_type = part.get_content_type()
                    content_disposition = str(part.get("Content-Disposition"))
                    
                    if content_type == "text/plain" and "attachment" not in content_disposition:
                        payload = part.get_payload(decode=True)
                        if payload:
                            body_text += payload.decode(part.get_content_charset() or 'utf-8', errors='ignore') + "\n"
                    elif content_type == "text/html" and not body_text and "attachment" not in content_disposition:
                        payload = part.get_payload(decode=True)
                        if payload:
                            html_str = payload.decode(part.get_content_charset() or 'utf-8', errors='ignore')
                            soup = BeautifulSoup(html_str, 'html.parser')
                            body_text += soup.get_text() + "\n"
            else:
                payload = msg.get_payload(decode=True)
                if payload:
                    body_text = payload.decode(msg.get_content_charset() or 'utf-8', errors='ignore')
            
            clean_body = body_text.strip()
            header_block = f"From: {headers['from']}\nTo: {headers['to']}\nSubject: {headers['subject']}\nDate: {headers['date']}\n\n"
            full_text = header_block + clean_body
            
            return ParsedDocument(
                source_name=source_name,
                mime_type="message/rfc822",
                sections=[ParsedSection(text=full_text, metadata={"section": "Email Body", "headers": headers, "source": source_name})],
                success=True
            )
        except Exception as e:
            return ParsedDocument(
                source_name=source_name,
                mime_type="message/rfc822",
                sections=[],
                success=False,
                error_message=f"EML parser error: {str(e)}"
            )

    @staticmethod
    def parse_txt(file_path: str, source_name: str) -> ParsedDocument:
        try:
            encoding = 'utf-8'
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
            except UnicodeDecodeError:
                encoding = 'latin-1'
                with open(file_path, 'r', encoding=encoding) as f:
                    content = f.read()
            
            clean_content = content.strip()
            if not clean_content:
                return ParsedDocument(
                    source_name=source_name,
                    mime_type="text/plain",
                    sections=[],
                    success=False,
                    error_message="Text file is empty."
                )

            return ParsedDocument(
                source_name=source_name,
                mime_type="text/plain",
                sections=[ParsedSection(text=clean_content, metadata={"section": "Content", "source": source_name})],
                success=True
            )
        except Exception as e:
            return ParsedDocument(
                source_name=source_name,
                mime_type="text/plain",
                sections=[],
                success=False,
                error_message=f"TXT parser error: {str(e)}"
            )

    @staticmethod
    def parse_image_ocr(file_path: str, source_name: str) -> ParsedDocument:
        if pytesseract is None:
            return ParsedDocument(
                source_name=source_name,
                mime_type="image/unknown",
                sections=[],
                success=False,
                error_message="pytesseract is not installed or available on this system."
            )
        try:
            img = Image.open(file_path)
            # Run OCR
            text = pytesseract.image_to_string(img)
            clean_text = text.strip()
            
            if not clean_text:
                return ParsedDocument(
                    source_name=source_name,
                    mime_type="image/ocr",
                    sections=[],
                    success=False,
                    error_message="OCR returned no text contents for this image."
                )

            return ParsedDocument(
                source_name=source_name,
                mime_type="image/ocr",
                sections=[ParsedSection(text=clean_text, metadata={"section": "OCR Results", "source": source_name})],
                success=True
            )
        except Exception as e:
            return ParsedDocument(
                source_name=source_name,
                mime_type="image/ocr",
                sections=[],
                success=False,
                error_message=f"OCR engine error: {str(e)}"
            )

    @classmethod
    def parse(cls, file_path: str, source_name: str) -> ParsedDocument:
        """
        Routes the file to the appropriate parser after sniffing the MIME type.
        Ensures that if any parsing operation fails, the system catches it gracefully.
        """
        mime_type = detect_mime_type(file_path)
        logger.info(f"Detected MIME type: {mime_type} for {source_name}")

        try:
            if mime_type == "application/pdf":
                return cls.parse_pdf(file_path, source_name)
            elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                return cls.parse_docx(file_path, source_name)
            elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                return cls.parse_pptx(file_path, source_name)
            elif mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                return cls.parse_xlsx(file_path, source_name)
            elif mime_type == "text/csv":
                return cls.parse_csv(file_path, source_name)
            elif mime_type == "text/html":
                return cls.parse_html(file_path, source_name)
            elif mime_type == "text/markdown":
                return cls.parse_markdown(file_path, source_name)
            elif mime_type == "message/rfc822":
                return cls.parse_eml(file_path, source_name)
            elif mime_type == "text/plain":
                return cls.parse_txt(file_path, source_name)
            elif mime_type.startswith("image/"):
                return cls.parse_image_ocr(file_path, source_name)
            else:
                return ParsedDocument(
                    source_name=source_name,
                    mime_type=mime_type,
                    sections=[],
                    success=False,
                    error_message=f"Unsupported format or MIME type: {mime_type}"
                )
        except Exception as e:
            # Catches unexpected parser crashes to satisfy the graceful failure requirement
            return ParsedDocument(
                source_name=source_name,
                mime_type=mime_type,
                sections=[],
                success=False,
                error_message=f"Fatal parsing error: {str(e)}"
            )
