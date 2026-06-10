import os
import csv
import zipfile
import pypdf

# Try importing reportlab for PDF generation, or fail gracefully
try:
    from reportlab.pdfgen import canvas
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False

try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import pptx
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False


def create_test_files():
    dir_name = "test_files"
    os.makedirs(dir_name, exist_ok=True)
    print(f"Creating sample test files in '{dir_name}/'...")

    # 1. Plain Text Notes (notes.txt)
    with open(os.path.join(dir_name, "notes.txt"), "w", encoding="utf-8") as f:
        f.write("Office Security Protocol and Codes\n")
        f.write("==================================\n\n")
        f.write("- Office door code is 9081#.\n")
        f.write("- Alarm system code is 4421.\n")
        f.write("- Please make sure the door is locked when leaving after 8 PM.\n")
    print("✓ Created notes.txt")

    # 2. Markdown File (readme.md)
    with open(os.path.join(dir_name, "readme.md"), "w", encoding="utf-8") as f:
        f.write("# RAG System Local Run Guide\n\n")
        f.write("To build and execute this project locally, run:\n")
        f.write("```bash\n")
        f.write("python run.py --port 8000\n")
        f.write("```\n")
        f.write("Ensure your environment variables are configured in `.env` first.\n")
    print("✓ Created readme.md")

    # 3. CSV File (sales.csv)
    with open(os.path.join(dir_name, "sales.csv"), "w", encoding="utf-8", newline="") as f:
        writer = csv.writer(f)
        writer.writerow(["Region", "Manager", "Target_Met", "Q1_Sales"])
        writer.writerow(["North", "Jane Doe", "Yes", "120000"])
        writer.writerow(["South", "John Smith", "No", "85000"])
        writer.writerow(["East", "Alice Williams", "Yes", "150000"])
    print("✓ Created sales.csv")

    # 4. HTML Homepage (homepage.html)
    with open(os.path.join(dir_name, "homepage.html"), "w", encoding="utf-8") as f:
        f.write("<!DOCTYPE html>\n<html>\n<head>\n<title>AURA Corporation Homepage</title>\n</head>\n")
        f.write("<body>\n")
        f.write("<h1>Welcome to AURA Corporation</h1>\n")
        f.write("<p>Our core mission statement is: <strong>Empowering document intelligence.</strong></p>\n")
        f.write("<p>Contact us at contact@auracorp.example.com for partnerships.</p>\n")
        f.write("</body>\n</html>\n")
    print("✓ Created homepage.html")

    # 5. EML Ticket (support_ticket.eml)
    with open(os.path.join(dir_name, "support_ticket.eml"), "w", encoding="utf-8") as f:
        f.write("From: support-system@auracorp.example.com\n")
        f.write("To: engineering-team@auracorp.example.com\n")
        f.write("Subject: Critical login failure reported by client X\n")
        f.write("Date: Mon, 08 Jun 2026 14:22:10 -0400\n")
        f.write("Content-Type: text/plain; charset=\"utf-8\"\n\n")
        f.write("Dear team,\n\n")
        f.write("We received an urgent support ticket. The customer was unable to log in because port 443 was blocked by their firewall.\n")
        f.write("Please advise them on how to unblock port 443 in their network configuration.\n\n")
        f.write("Best,\nSupport Team\n")
    print("✓ Created support_ticket.eml")

    # 6. Word Document (engineering_handbook.docx)
    if DOCX_AVAILABLE:
        doc = docx.Document()
        doc.add_heading("AURA Engineering Handbook", level=0)
        doc.add_heading("Git Workflow Guidelines", level=1)
        doc.add_paragraph("All developers must adhere to the standardized git naming conventions.")
        doc.add_paragraph("Git branch names must follow feature/feature-name or bugfix/bug-name.")
        doc.save(os.path.join(dir_name, "engineering_handbook.docx"))
        print("✓ Created engineering_handbook.docx")
    else:
        print("✗ Skipped docx (python-docx not installed yet)")

    # 7. PowerPoint Slide (product_roadmap.pptx)
    if PPTX_AVAILABLE:
        prs = pptx.Presentation()
        # Slide 1
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        slide1.shapes.title.text = "Product Roadmap 2026"
        slide1.placeholders[1].text = "Roadmap outline for next major versions."
        
        # Slide 2
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Version 2.0 Details"
        slide2.placeholders[1].text = "Version 2.0 is scheduled for release in Q4 2026.\nIt features document analysis and multi-modal citations."
        
        prs.save(os.path.join(dir_name, "product_roadmap.pptx"))
        print("✓ Created product_roadmap.pptx")
    else:
        print("✗ Skipped pptx (python-pptx not installed yet)")

    # 8. Excel spreadsheet (quarterly_earnings.xlsx)
    if OPENPYXL_AVAILABLE:
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Q3_Earnings"
        ws.append(["Category", "Metric", "Amount"])
        ws.append(["Finance", "Q3 Revenue", "$4.2M"])
        ws.append(["Finance", "Net Profit", "$850k"])
        ws.append(["Finance", "Gross Margin", "78%"])
        wb.save(os.path.join(dir_name, "quarterly_earnings.xlsx"))
        print("✓ Created quarterly_earnings.xlsx")
    else:
        print("✗ Skipped xlsx (openpyxl not installed yet)")

    # 9. PDF document (company_policy.pdf)
    if REPORTLAB_AVAILABLE:
        c = canvas.Canvas(os.path.join(dir_name, "company_policy.pdf"))
        c.setFont("Helvetica-Bold", 16)
        c.drawString(50, 750, "AURA Corp Company Policy Handbook")
        
        c.setFont("Helvetica-Bold", 12)
        c.drawString(50, 700, "Section 1: Remote Work Allowance")
        c.setFont("Helvetica", 10)
        c.drawString(50, 680, "Employees are eligible for home office expenses configuration.")
        c.drawString(50, 665, "The remote work allowance is $500 per year.")
        
        c.setFont("Helvetica-Bold", 12)
        c.drawString(50, 600, "Section 2: Health Insurance")
        c.setFont("Helvetica", 10)
        c.drawString(50, 580, "Medical benefits cover 100% of standard premiums.")
        c.save()
        print("✓ Created company_policy.pdf")
    else:
        print("✗ Skipped pdf (reportlab not installed yet)")

    # 10. Corrupt File (corrupt.pdf) - intentionally bad bytes to test parser robustness
    with open(os.path.join(dir_name, "corrupt.pdf"), "wb") as f:
        f.write(b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\nThis is a corrupted PDF file. We have random bytes \xff\xfe\xfd\x00\x01\x02 which cannot be parsed by any pdf reader.")
    print("✓ Created corrupt.pdf")

    # 11. Password protected PDF (password.pdf)
    try:
        writer = pypdf.PdfWriter()
        writer.add_blank_page(width=300, height=300)
        writer.encrypt("securepwd123")
        with open(os.path.join(dir_name, "password.pdf"), "wb") as f:
            writer.write(f)
        print("✓ Created password.pdf")
    except Exception as e:
        print(f"✗ Failed to create password.pdf: {e}")


if __name__ == "__main__":
    create_test_files()
